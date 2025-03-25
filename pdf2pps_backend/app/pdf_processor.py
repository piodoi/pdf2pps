import fitz  # PyMuPDF
from typing import List, Dict, Any, Optional
import os
import tempfile
import gc
from pptx import Presentation
from pptx.util import Inches, Pt
from .models import PresentationSlide, PresentationResponse


class PDFProcessor:
    """Class to handle PDF processing and presentation generation."""
    
    def __init__(self, model_name: str = "llama2"):
        """Initialize with the specified LLM model."""
        self.model_name = model_name
        
    def extract_text_from_pdf(self, file_path: str, max_pages: int = 20) -> str:
        """Extract text content from a PDF file with memory optimization."""
        text = ""
        try:
            # Open the document
            doc = fitz.open(file_path)
            
            # Limit number of pages to process to avoid memory issues
            page_count = min(doc.page_count, max_pages)
            
            # Process pages one by one to reduce memory usage
            for page_num in range(page_count):
                page = doc.load_page(page_num)
                text += page.get_text()
                # Explicitly delete page object to free memory
                page = None
                # Force garbage collection
                gc.collect()
                
            # Close the document
            doc.close()
            doc = None
            gc.collect()
            
            return text
        except Exception as e:
            raise Exception(f"Error extracting text from PDF: {str(e)}")
    
    def extract_key_points(self, text: str, max_points: int = 6) -> List[PresentationSlide]:
        """Extract key points using rule-based approach with memory optimization."""
        try:
            # Use only rule-based approach for deployment to avoid OOM issues
            import re
            
            # Limit text size to reduce memory usage
            text = text[:20000]
            
            # Simple rule-based extraction
            # Split text into paragraphs more efficiently
            paragraphs = text.split('\n\n')
            sentences = []
            
            # Process paragraphs in chunks to reduce memory usage
            for para in paragraphs:
                # Split into sentences
                para_sentences = re.split(r'(?<=[.!?])\s+', para)
                # Only keep non-empty sentences
                for s in para_sentences:
                    s_stripped = s.strip()
                    if s_stripped:
                        sentences.append(s_stripped)
            
            # Extract bullet points and numbered lists
            bullet_points = []
            for line in text.split('\n'):
                if re.match(r'^\s*(\d+\.|•|\*|-)\s+', line):
                    content = re.sub(r'^\s*(\d+\.|•|\*|-)\s+', '', line).strip()
                    if content:
                        bullet_points.append(content)
            
            # Create slides with memory optimization
            slides = []
            
            # Title slide
            slides.append(PresentationSlide(
                title="Document Summary",
                content=["Generated from PDF document", "Key points extracted"]
            ))
            
            # Introduction slide with first few sentences
            if sentences:
                intro_content = sentences[:min(3, len(sentences))]
                slides.append(PresentationSlide(
                    title="Introduction",
                    content=intro_content
                ))
            
            # Bullet points slide
            if bullet_points:
                slides.append(PresentationSlide(
                    title="Key Points",
                    content=bullet_points[:min(5, len(bullet_points))]
                ))
            
            # Content slides from remaining text
            remaining_sentences = sentences[min(3, len(sentences)):]
            
            # Group sentences into slides (3-4 per slide)
            for i in range(0, len(remaining_sentences), 4):
                group = remaining_sentences[i:i+4]
                if group:
                    slides.append(PresentationSlide(
                        title=f"Content {i//4 + 1}",
                        content=group
                    ))
                
                # Limit to max_points slides
                if len(slides) >= max_points:
                    break
            
            # Conclusion slide
            if len(slides) < max_points and sentences:
                slides.append(PresentationSlide(
                    title="Conclusion",
                    content=["End of document summary", 
                            "For more detailed information, please refer to the original document"]
                ))
            
            # Force garbage collection
            gc.collect()
            
            return slides
                
        except Exception as e:
            # Return error slide
            return [
                PresentationSlide(
                    title="Error Processing Document",
                    content=[f"An error occurred: {str(e)}", 
                             "Please try again or check if the PDF contains extractable text"]
                )
            ]
    
    def generate_pptx(self, slides: List[PresentationSlide], output_filename: str) -> str:
        """Generate a PPTX presentation from the slides with memory optimization."""
        prs = Presentation()
        
        # Create title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Document Summary"
        subtitle.text = "Generated by PDF2PPS"
        
        # Create content slides
        for slide_data in slides:
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data.title
            
            # Add content as bullet points
            content = slide.placeholders[1]
            tf = content.text_frame
            
            for point in slide_data.content:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
        
        # Save the presentation
        prs.save(output_filename)
        
        # Force garbage collection
        gc.collect()
        
        return output_filename
    
    def process_pdf(self, pdf_path: str, output_dir: str = "") -> str:
        """Process a PDF file and generate a presentation with memory optimization."""
        if not output_dir:
            output_dir = tempfile.gettempdir()
            
        # Extract text from PDF
        text = self.extract_text_from_pdf(pdf_path)
        
        # Extract key points using rule-based approach
        slides = self.extract_key_points(text)
        
        # Generate filename for the presentation
        base_filename = os.path.splitext(os.path.basename(pdf_path))[0]
        output_path = os.path.join(output_dir, f"{base_filename}_presentation.pptx")
        
        # Generate the presentation
        self.generate_pptx(slides, output_path)
        
        return output_path
